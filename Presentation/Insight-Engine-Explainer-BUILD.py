#!/usr/bin/env python3
# Self-contained generator for Insight-Engine-Explainer.pptx (27 slides, Anthropic palette).
# Requires: pip install python-pptx --break-system-packages
# Run: python3 Insight-Engine-Explainer-BUILD.py  -> writes the .pptx next to this script.
# This is the consolidated final build (was build_deck_v2 -> patch_v3..v8 in an ephemeral scratchpad).

#!/usr/bin/env python3
"""Build the Insight Engine chaptered explainer deck (v2, 27 slides)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# palette
# Anthropic brand palette
IVORY=RGBColor(0xFA,0xF9,0xF5)      # page background (light)
NAVY_BG=RGBColor(0x14,0x14,0x13)    # dark slides (Anthropic Dark)
INK=RGBColor(0x14,0x14,0x13)        # primary text
BODY=RGBColor(0x4A,0x45,0x3D)       # body / secondary text (warm dark grey)
MUTE=RGBColor(0x8A,0x85,0x78)       # captions / muted (warm grey)
BLUE=RGBColor(0x4E,0x79,0xA6)       # Anthropic secondary blue (darkened for text)
GREEN=RGBColor(0x5F,0x72,0x49)      # Anthropic tertiary green (darkened for text)
CORAL=RGBColor(0xD9,0x77,0x57)      # Anthropic primary accent (coral)
AMBER=CORAL                         # remap old amber -> coral
GOLD=RGBColor(0xE8,0x90,0x78)       # light coral (accent on dark)
DGOLD=RGBColor(0xA9,0x6A,0x45)      # warm brown for [V3]
DAMBER=RGBColor(0x7A,0x3F,0x28)     # dark coral-brown text on coral tint
ICE=RGBColor(0xFA,0xF9,0xF5)        # light text on dark
ICE2=RGBColor(0xC9,0xC6,0xBA)       # secondary light on dark
WHITE=RGBColor(0xFF,0xFF,0xFF)
CARD=RGBColor(0xFF,0xFF,0xFF)       # white cards on ivory
BORD=RGBColor(0xE4,0xE1,0xD6)       # warm border
GTINT=RGBColor(0xEB,0xEE,0xE2)      # green tint
ATINT=RGBColor(0xF2,0xDB,0xCE)      # coral tint
F="Arial"
CARDGEO=[]; CURSLIDE=[0]

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]

def slide(bg=IVORY):
    s=prs.slides.add_slide(BLANK)
    CURSLIDE[0]+=1
    f=s.background.fill; f.solid(); f.fore_color.rgb=bg
    return s

def note(s,txt):
    s.notes_slide.notes_text_frame.text=txt

def tb(s,x,y,w,h,runs,size,color,bold=False,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,font=F,italic=False,sp=None):
    b=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=b.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    for m in ("left","right","top","bottom"): setattr(tf,f"margin_{m}",0)
    if isinstance(runs,str): runs=[runs]
    def apply(p,seg):
        run=p.add_run(); run.text=seg[0]; fnt=run.font
        fnt.size=Pt(seg[1] if len(seg)>1 and seg[1] else size)
        fnt.bold=seg[2] if len(seg)>2 and seg[2] is not None else bold
        fnt.color.rgb=seg[3] if len(seg)>3 and seg[3] else color
        fnt.name=seg[4] if len(seg)>4 and seg[4] else font
        fnt.italic=seg[5] if len(seg)>5 and seg[5] else italic
    for i,r in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align
        if sp is not None: p.space_after=Pt(sp)
        if isinstance(r,str):
            run=p.add_run(); run.text=r; fnt=run.font
            fnt.size=Pt(size); fnt.bold=bold; fnt.color.rgb=color; fnt.name=font; fnt.italic=italic
        elif isinstance(r,tuple):
            apply(p,r)
        elif isinstance(r,list):
            for seg in r: apply(p,seg)
    return b

def card(s,x,y,w,h,fill=CARD,border=BORD,line=True,round=True):
    CARDGEO.append({"slide":CURSLIDE[0],"x":x,"y":y,"w":w,"h":h})
    shp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if round else MSO_SHAPE.RECTANGLE,
                           Inches(x),Inches(y),Inches(w),Inches(h))
    try: shp.adjustments[0]=0.055
    except: pass
    shp.fill.solid(); shp.fill.fore_color.rgb=fill
    if line: shp.line.color.rgb=border; shp.line.width=Pt(1)
    else: shp.line.fill.background()
    shp.shadow.inherit=False
    return shp

def label(s,txt,color=CORAL):
    tb(s,0.72,0.36,11.9,0.4,txt,15,color,bold=True)
def headline(s,txt,color=INK,size=30):
    tb(s,0.70,0.74,11.93,1.0,txt,size,color,bold=True)
def banner(s,parts,y=5.75,h=1.15,fill=NAVY_BG,color=WHITE):
    card(s,0.70,y,11.93,h,fill=fill,border=fill,line=False)
    tb(s,1.10,y,11.13,h,parts,18,color,bold=True,anchor=MSO_ANCHOR.MIDDLE)

N=[]  # narration notes, filled per slide

# ---------- S1 TITLE ----------
s=slide()
tb(s,0.90,1.30,11.5,0.7,"Insight Engine",30,CORAL,bold=True)
tb(s,0.88,2.15,11.6,1.4,"Analysis defended with assurance.",54,INK,bold=True)
tb(s,0.92,4.15,10.9,1.7,[[("Insight",22,True,BODY),(" supplied by a Large Language Model;",22,False,BODY)],[("Engine",22,True,BODY),(" makes that insight verified, defensible, and audience-ready.",22,False,BODY)]],22,BODY,sp=6)
tb(s,0.92,6.5,11.2,0.4,"An open plugin for Claude Code & Cowork   ·   v0.1.17",16,MUTE)
note(s,"This is the Insight Engine — a tool for analysis you can actually defend. The idea behind it fits in one sentence: the language model supplies the insight; the engine makes that insight verified, graded, and honest about its own limits. It doesn't try to think better than the model. It makes the model show its work — and checks that work against the world.")

# ---------- S1b THE MAP (overview) ----------
s=slide()
label(s,"THE MAP")
headline(s,"One architecture — here is the whole of it.")
_ov=[("THE IDEA — two jobs, one rule","Reveal what's hidden  ·  Verify what's there  ·  the value calls stay yours",BLUE),
     ("THE PIPELINE — runs silently, except the gate","Triage → Provocation → Verify & grade → Disconfirm → Systems pass (conditional) → Deep-core routing → Grade-lock → Forced gate → Decision brief",GREEN),
     ("THE INTERFACE — five commands","analyse  ·  verify  ·  render  ·  track  ·  adjudicate",BLUE),
     ("THE MODEL LAYER — execution + assurance","A frontier model executes. For the highest stakes, an independent adjudication pass (Step 10) — opt-in, off by default.",GREEN)]
for _i,(_hd,_bd,_c) in enumerate(_ov):
    _col=_i%2; _row=_i//2; _x=0.70+_col*6.07; _y=1.80+_row*1.92
    card(s,_x,_y,5.80,1.74)
    tb(s,_x+0.30,_y+0.18,5.2,0.5,_hd,15,_c,bold=True)
    tb(s,_x+0.30,_y+0.70,5.2,0.95,_bd,13,BODY)
banner(s,["Architected around ICD 203 — the U.S. intelligence community's Analytic Standards.","Each part is explained in the slides ahead."],y=5.60,h=1.30)
note(s,"Before we walk it, here is the map.")

# ---------- S1c INSIGHT + ENGINE ----------
s=slide()
label(s,"THE NAME IS THE DESIGN")
headline(s,"Insight from the model. Assurance from the engine.")
card(s,0.70,1.95,5.95,3.30,fill=CARD,border=BLUE)
tb(s,1.0,2.2,5.4,0.6,"INSIGHT · a strong LLM",20,BLUE,bold=True)
tb(s,1.0,3.05,5.4,2.1,"The language model supplies the raw material — the reasoning, the reading, the first draft, the searches. All the intelligence comes from the model; the engine assumes it.",17,BODY)
card(s,6.85,1.95,5.78,3.30,fill=GTINT,border=GREEN)
tb(s,7.15,2.2,5.2,0.6,"ENGINE · verify · assure · adjudicate",20,GREEN,bold=True)
tb(s,7.15,3.05,5.2,2.1,[("• Verify — check what can be checked against the world",16,False,INK),("• Assure — grade how strongly each claim stands, and lock it",16,False,INK),("• Adjudicate — for the highest stakes, an independent model checks the result",16,False,INK)],16,INK,sp=8)
banner(s,"The model provides the insight; the engine makes it verified, defensible, and — when it matters most — independently checked.",y=5.60,h=1.1)
note(s,"Insight from the model; verification, assurance, and adjudication from the engine.")

# ---------- S2 PROBLEM ----------
s=slide()
label(s,"THE PROBLEM"); headline(s,"A confident guess is not verified evidence.")
cols=[("The important things hide","What matters most sits where attention doesn't go."),
      ("The evidence is mixed, or interested","Sources disagree — and some benefit from the answer."),
      ("The deciding questions are value judgments","No analysis settles what you should value.")]
for i,(hd,bd) in enumerate(cols):
    x=0.70+i*4.07; card(s,x,1.95,3.80,2.75)
    tb(s,x+0.3,2.2,3.2,1.0,hd,19,BLUE,bold=True)
    tb(s,x+0.3,3.35,3.2,1.25,bd,17,BODY)
banner(s,[("A strong model gives you a fluent answer.  ",18,True,WHITE),("Fluent isn't the same as verified — or defensible.",18,True,ICE)],y=5.05,h=1.4)
note(s,"A confident guess is not verified evidence. High-stakes, complex, or wicked problems — the ones with no clean solution — tend to share three features. The things that matter most hide — they sit where attention doesn't naturally go. The evidence is mixed, or it comes from people with a stake in the answer. And the questions that actually decide the outcome are value judgments, which no amount of analysis can settle for you. Hand a problem like that to a strong model and you get a fluent, confident answer. But fluent isn't the same as verified — and it certainly isn't the same as defensible.")

# ---------- S3 IDEA ----------
s=slide()
label(s,"THE IDEA"); headline(s,"Two jobs: reveal what's hidden, and verify what's there.")
card(s,0.70,1.95,5.95,3.25,fill=CARD,border=BLUE)
tb(s,1.0,2.2,5.4,0.6,"1 · Reveal what's hidden",20,BLUE,bold=True)
tb(s,1.0,3.0,5.4,2.1,[("• the voice that's absent",16,False,INK),("• what's been normalised until no one questions it",16,False,INK),
    ("• the thing left conspicuously out",16,False,INK),("• some missed — some kept out of view on purpose",16,False,INK)],16,INK,sp=6)
card(s,6.85,1.95,5.78,3.25,fill=GTINT,border=GREEN)
tb(s,7.15,2.2,5.2,0.6,"2 · Verify what's there",20,GREEN,bold=True)
tb(s,7.15,3.0,5.2,2.1,[("• check what can be checked",16,False,INK),("• grade how strongly each claim stands",16,False,INK),
    ("• try to break its own findings",16,False,INK)],16,INK,sp=6)
banner(s,"Seeing and checking — then it makes you own the value calls, instead of quietly making them for you.",y=5.5,h=1.1)
note(s,"So the engine's move isn't to be smarter — it's to be more honest, and it has two jobs, not one. The first is to reveal what's hidden: the voice that's absent, the fact that's been normalised until no one questions it, the thing left conspicuously out — some of it simply missed, some kept out of view on purpose. The second is to verify what's there: check what can be checked, grade how strongly each checkable claim stands, and try to break its own findings. Seeing and checking — and then it makes you own the value calls, instead of quietly making them for you. Everything that follows serves those two jobs.")

# ---------- S4 HOW IT WORKS ----------
s=slide()
label(s,"HOW IT WORKS"); headline(s,"What happens when you hand it a problem")
A="/insight-engine:analyse"
GY=1.50; GP=1.38; CW=3.75; CH=1.26; COLP=4.05
cells=[
 ("STEP 1","Triage","size the effort first",BLUE,A,CARD,BORD,MUTE,BODY),
 ("STEP 2","Scoping","optional clarifying Qs",BLUE,A,CARD,BORD,MUTE,BODY),
 ("STEP 3","Provocation","surface what hides",BLUE,A,CARD,BORD,MUTE,BODY),
 ("STEP 4","Verification","grade every claim",GREEN,"/insight-engine:verify",CARD,BORD,MUTE,BODY),
 ("STEP 5","Systems pass","conditional",BLUE,A,CARD,BORD,MUTE,BODY),
 ("STEP 6","Deep-core","set aside un-answerable",BLUE,A,CARD,BORD,MUTE,BODY),
 ("STEP 7","Grade-lock","fix + stress-test",GREEN,A,CARD,BORD,MUTE,BODY),
 ("STEP 8","THE GATE","your judgment",AMBER,A,CARD,BORD,MUTE,BODY),
 ("STEP 9","Decision brief","the call · what flips it",BLUE,A,CARD,BORD,MUTE,BODY),
 ("STEP 10","Adjudication","opt-in · 2nd-model check",DAMBER,"/insight-engine:adjudicate",ATINT,CORAL,DAMBER,DAMBER),
 ("STEP 11","Render","re-voice for a reader",DAMBER,"/insight-engine:render",ATINT,CORAL,DAMBER,DAMBER),
 ("STEP 12","Track","keep it live as facts arrive",DAMBER,"/insight-engine:track",ATINT,CORAL,DAMBER,DAMBER),
]
for i,(st,hd,bd,c,cmd,fl,br,cs,cb) in enumerate(cells):
    col=i%3; row=i//3; x=0.70+col*COLP; y=GY+row*GP
    card(s,x,y,CW,CH,fill=fl,border=br)
    tb(s,x+0.2,y+0.07,3.4,0.24,st,10,cs,bold=True)
    tb(s,x+0.2,y+0.30,3.4,0.40,hd,15,c,bold=True)
    tb(s,x+0.2,y+0.70,3.4,0.26,bd,12,cb)
    tb(s,x+0.2,y+0.98,3.4,0.22,cmd,8.5,cs,font="Consolas")
    if col<2: tb(s,x+3.77,y,0.26,CH,"‣",18,cs,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
note(s,"S4 overview note")

# ---------- S5 TRIAGE ----------
s=slide()
label(s,"STEP 1 · PROPORTIONALITY TRIAGE"); headline(s,"First, decide how hard to look.")
card(s,0.70,1.95,7.35,3.55)
tb(s,1.00,2.2,6.8,0.5,"Runs the FULL method — the default",19,BLUE,bold=True)
tb(s,1.00,2.9,6.8,2.5,[("Any one of these is enough:",15,False,BODY,F,True),
    ("• irreversible",16,False,INK),("• someone bound who isn't in the room",16,False,INK),
    ("• real money or dependency",16,False,INK),("• legal or regulatory exposure",16,False,INK),
    ("• contested values",16,False,INK),("• a party who benefits from light scrutiny",16,False,INK)],16,INK,sp=4)
card(s,8.30,1.95,4.33,3.55,fill=CARD)
tb(s,8.60,2.2,3.75,0.5,"A short answer — rare",19,BODY,bold=True)
tb(s,8.60,2.9,3.75,2.4,"Only when the problem is unambiguously routine, reversible, and binds nobody. Under any doubt, it runs the full method.",17,BODY)
banner(s,[("A little over-analysis is cheap.  ",18,True,WHITE),("Treating a serious problem as routine is not.",18,True,ICE)])
note(s,"The first thing it does is decide how hard to look. Not every question deserves the full machine — but that call is deliberately fail-safe. Depth is the default. It only scales down when a problem is unambiguously routine, reversible, and binds nobody. If there's any doubt at all — anything irreversible, anyone affected who isn't in the room, real money, legal exposure, or a party who benefits from light scrutiny — it runs the full method. A little over-analysis is cheap. Treating a serious problem as routine is not.")

# ---------- S5b SCOPING (Step 2) ----------
s=slide()
label(s,"STEP 2 · OPTIONAL SCOPING")
headline(s,"Five questions you can answer — or skip.")
card(s,0.70,1.95,11.93,3.35)
tb(s,1.05,2.18,11.2,3.0,[("Before the analysis, the engine offers up to five clarifying questions — one thing each, all skippable:",16,False,BODY),
   ("1 · The trouble — what is going wrong, stated as the symptom not the diagnosis, and why now.",15,False,INK),
   ("2 · The owner — whose problem this is, who acts on the answer, and who else will read it.",15,False,INK),
   ("3 · Success & stakes — what a good outcome looks like, what counts as a serious cost, whose costs matter most.",15,False,INK),
   ("4 · Constraints & scope — what is fixed, and what is explicitly in and out (scope is where the largest hiddens hide).",15,False,INK),
   ("5 · Party-held material — documents or figures only you hold; upload them.",15,False,INK)],15,INK,sp=7)
banner(s,["Skip any or all — nothing blocks.","Whatever you don't answer, the engine fills with a stated assumption and marks it as one."],y=5.55,h=1.30)
note(s,"S5b scoping note")

# ---------- S6 PROVOCATION ----------
s=slide()
label(s,"STEP 3 · THE PROVOCATION PASS"); headline(s,"See what's hidden — before you verify.")
probes=[("Whose voice is absent?","who isn't in the room"),("Who benefits from the framing?","cui bono"),
        ("What's normalised","treated as natural — but shouldn't be"),("What's conspicuously MISSING?","the document that isn't there; the un-asked question")]
for i,(hd,bd) in enumerate(probes):
    col=i%2; row=i//2; x=0.70+col*6.07; y=1.85+row*1.42
    hi = i==3
    card(s,x,y,5.80,1.26,fill=ATINT if hi else CARD, border=GOLD if hi else BORD)
    tb(s,x+0.28,y+0.15,5.3,0.55,hd,18,DAMBER if hi else BLUE,bold=True)
    tb(s,x+0.28,y+0.70,5.3,0.45,bd,14,BODY)
tb(s,0.70,4.82,11.9,0.75,"Some is simply overlooked; some is kept out of view on purpose. Either way — make the invisible visible.",17,INK,bold=True)
banner(s,"Then it steelmans your strongest findings and names everyone the outcome materially touches — including those not in the room.",y=5.75,h=1.15)
note(s,"Then, before it verifies anything, it goes hunting for your blind spots — the things a fluent first reading glides straight past. This is the provocation pass — a disciplined set of probes, a couple always run and the rest chosen to fit the case, aimed at what the page hides and at what it quietly smuggles in. Whose voice is absent? Who benefits from the way the problem's been framed? What's being treated as normal that shouldn't be? And hardest of all — what's conspicuously missing: the document that isn't there, the question nobody asked? Some of what matters is simply overlooked; some is kept out of view by someone it would inconvenience. Either way, the engine's job here is to make the invisible visible — then it steelmans your strongest findings and names everyone the outcome materially touches, including those not in the room — so you decide with the whole picture, not the convenient part of it.")

# ---------- S7 VERIFY & GRADE-LOCK ----------
s=slide()
label(s,"STEP 4 · VERIFY & GRADE",GREEN); headline(s,"Every claim — graded, and locked.")
card(s,0.70,1.75,5.80,1.15)
tb(s,1.0,1.75,1.5,1.15,"[V]",34,GREEN,bold=True,anchor=MSO_ANCHOR.MIDDLE)
tb(s,2.35,1.75,3.9,1.15,"verified — independently corroborated",17,INK,anchor=MSO_ANCHOR.MIDDLE)
card(s,6.83,1.75,5.80,1.15)
tb(s,7.13,1.75,1.5,1.15,"[N]",34,AMBER,bold=True,anchor=MSO_ANCHOR.MIDDLE)
tb(s,8.48,1.75,3.9,1.15,"not — rests on an interested party",17,INK,anchor=MSO_ANCHOR.MIDDLE)
tiers=[("[V1]","primary",GREEN),("[V2]","secondary",BLUE),("[V3]","weak / contested",DGOLD)]
for i,(t,d,c) in enumerate(tiers):
    x=0.70+i*4.07; card(s,x,3.10,3.80,0.85)
    tb(s,x+0.25,3.10,3.4,0.85,[(t+"  ",18,True,c),(d,16,False,BODY)],16,BODY,anchor=MSO_ANCHOR.MIDDLE)
card(s,0.70,4.20,11.93,1.15,fill=GTINT,border=GREEN)
tb(s,1.05,4.20,11.3,1.15,[("A tier marks how CLOSE a source is — not how much to trust it.  ",17,True,INK),
    ("A party's own document doesn't earn [V1] on a claim it has a stake in.",17,False,BODY)],17,INK,anchor=MSO_ANCHOR.MIDDLE)
banner(s,["Set here, and held — never changed downstream, only re-expressed.","The formal lock comes at Step 7."],y=5.60,h=1.30)
note(s,"Now it verifies. Every claim the conclusion leans on gets sorted into one of three: a public fact that can be checked, a private document only some party holds, or a value that no evidence can settle. The checkable ones are searched and graded. V, for verified, means independently corroborated. N, for not, means not independently verified — including when a claim rests only on an interested party's own account. And every V carries a strength tier — V1 primary, V2 secondary, V3 weak or contested. Here's the part the engine is careful about: a tier marks how close a source is, not how far you should trust it. A company's own incident report is a primary document — but on a claim that company has a stake in, its own paperwork does not earn a V1. That grade has to come from something independent. And once verification is complete — including the disconfirmation pass you'll see next — the grades lock: nothing downstream can change them, only restate them. That's what lets the analysis hold up after the fact.")

# ---------- S7b HOW A CLAIM GETS GRADED ----------
s=slide()
label(s,"STEP 4 · HOW A CLAIM GETS GRADED",GREEN); headline(s,"Route, search, tier by independence, lock.")
routers=[("Public fact","searched against independent sources",GTINT,GREEN),
         ("Party-held","parked — name the doc that would settle it",CARD,BLUE),
         ("Judgement call","a value / framing / forecast no search settles — routed to you at the gate, never graded",CARD,AMBER)]
for i,(hd,bd,fill,c) in enumerate(routers):
    x=0.70+i*4.07; card(s,x,1.72,3.80,1.30,fill=fill,border=c if fill!=CARD else BORD)
    tb(s,x+0.24,1.88,3.35,0.5,hd,17,c,bold=True)
    tb(s,x+0.24,2.32,3.35,0.80,bd,12,BODY)
card(s,0.70,3.28,7.35,2.5)
gtiers=[("[V1]","primary — independent of any interested party",GREEN),
        ("[V2]","reputable reporting resting on those",BLUE),
        ("[V3]","weak but independent, or credible sources disagree",DGOLD),
        ("[N]","nothing independent behind it, or an interested party's own account",AMBER)]
for i,(t,d,c) in enumerate(gtiers):
    y=3.48+i*0.57
    tb(s,1.0,y,1.15,0.5,t,17,c,bold=True)
    tb(s,2.2,y,5.65,0.55,d,14,INK,anchor=MSO_ANCHOR.MIDDLE)
card(s,8.30,3.28,4.33,2.5,fill=GTINT,border=GREEN)
tb(s,8.60,3.55,3.75,2.0,[("Graded by independence.",16,True,GREEN),
    ("A party's own document doesn't earn V1 on a claim it has a stake in.",15,False,BODY)],15,BODY,sp=7)
banner(s,"Then each claim is disconfirmed — searched at its own negation, which can downgrade — and only then does it stand (formal lock at Step 7).",y=5.95,h=1.05)
note(s,"So how does a load-bearing claim get its grade? First it's routed. Is it a public fact you can check, a private document some party holds, or a judgement call — a value, a framing, a forecast — that no search can settle? The judgement call is handed to you; the private document is parked with a note of exactly what would settle it. A public fact gets searched — against independent current sources. If it's backed by a primary source independent of any interested party — the regulation itself, the official statistic, the ruling — that's V1. Reputable reporting resting on those is V2. If independent backing exists but it's weak, or credible sources genuinely disagree, it's V3. If the search turns up nothing independent behind it — including when all it rests on is an interested party's own account — it's N. Then each of those claims is disconfirmed — searched at its own negation, which can pull a grade down — and only then does it lock. That's the machine: route, search, tier by strength and independence, try to break it, lock.")

# ---------- S8 DISCONFIRMATION ----------
s=slide()
label(s,"STEP 4 · DISCONFIRMATION",GREEN); headline(s,"It tries to prove itself wrong.")
tb(s,0.70,1.8,11.9,0.6,"For every claim the conclusion rests on, it runs a second search — aimed at the opposite.",18,INK)
flow=[("A load-bearing claim",BLUE),("Search its negation",AMBER),("Survives, or downgrades",GREEN)]
for i,(t,c) in enumerate(flow):
    x=0.70+i*4.27; card(s,x,2.55,3.70,1.35)
    tb(s,x+0.2,2.55,3.3,1.35,t,18,c,bold=True,anchor=MSO_ANCHOR.MIDDLE,align=PP_ALIGN.CENTER)
    if i<2: tb(s,x+3.72,2.55,0.55,1.35,"→",30,MUTE,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
card(s,0.70,4.15,11.93,1.15,fill=GTINT,border=GREEN)
tb(s,1.05,4.15,11.3,1.15,[("“Survived” needs two independent sources.  ",17,True,INK),
    ("A single origin echoed around is flagged “unopposed” — not survived.",17,False,BODY)],17,INK,anchor=MSO_ANCHOR.MIDDLE)
banner(s,"The direct counter to confirmation bias — the failure mode any single-pass search invites.",y=5.6,h=1.1)
note(s,"Grading alone isn't enough, because a single confirming search is exactly how a wrong fact survives. So for every claim the conclusion really rests on, the engine runs a second search — aimed at the opposite. It actively hunts for the evidence that would break the claim. If credible counter-evidence of comparable weight turns up, the claim is downgraded and the dispute is recorded. If a genuine search for the negation finds nothing, the claim is marked as having survived — but only when the support comes from at least two independent sources. If it all traces back to a single origin echoed around, that's flagged as unopposed, not survived. This is the direct counter to confirmation bias — the failure mode any single-pass search quietly invites.")

# ---------- S9 SYSTEMS PASS ----------
s=slide()
label(s,"STEP 5 · SYSTEMS PASS · CONDITIONAL"); headline(s,"When the problem feeds back on itself.")
card(s,0.70,1.95,5.95,3.35)
tb(s,1.0,2.2,5.4,0.5,"Runs only when 2+ signatures appear",18,BLUE,bold=True)
tb(s,1.0,2.95,5.4,2.2,[("• feedback loops",16,False,INK),("• delays between cause and effect",16,False,INK),
    ("• things that accumulate",16,False,INK),("• many interacting actors",16,False,INK),
    ("• a possible tipping point",16,False,INK),("• self-fulfilling expectations",16,False,INK)],16,INK,sp=4)
card(s,6.85,1.95,5.78,3.35,fill=CARD)
tb(s,7.15,2.2,5.2,0.5,"A graded hypothesis — not a simulation",18,GREEN,bold=True)
tb(s,7.15,2.95,5.2,2.2,"It will tell you a loop EXISTS. It will not tell you, as fact, which loop WINS — that routes back to you. Structure can be verified; whether or when it tips cannot.",17,BODY)
banner(s,["It states whose framing the map takes, and who's affected but missing.","A tidy map that hides its blind spot is false precision."],y=5.55,h=1.30)
note(s,"Some problems aren't a list of facts — they're a structure that feeds back on itself. When a problem shows more than one signature of real feedback — loops, delays, things that accumulate, a possible tipping point — the engine runs a conditional systems pass. It maps the problem as causal links, marks the reinforcing and balancing loops, and names where a small push would move the most. But it holds a hard line: the map is a hypothesis to investigate, not a simulation to trust. It will tell you a loop exists; it will not tell you, as established fact, which loop wins — that routes back to you. And it states plainly whose framing the map takes and who is affected but missing from it. A tidy map that hides its own blind spot is false precision — and the engine treats that as a failure, not thoroughness.")

# ---------- S9b DEEP-CORE ROUTING (Step 3) ----------
s=slide()
label(s,"STEP 6 · DEEP-CORE ROUTING")
headline(s,"Separate out what no evidence can settle.")
card(s,0.70,1.95,7.35,3.35)
tb(s,1.0,2.2,6.8,0.5,"Onto an explicit OPEN list",18,BLUE,bold=True)
tb(s,1.0,2.95,6.8,2.25,[("• how the decision is framed in the first place",16,False,INK),("• legitimacy — who is bound without a say",16,False,INK),("• opportunity cost, and alternatives no one modelled",16,False,INK),("• the value trade-offs",16,False,INK)],16,INK,sp=7)
card(s,8.30,1.95,4.33,3.35,fill=ATINT,border=CORAL)
tb(s,8.60,2.2,3.75,0.5,"Three rules",18,DAMBER,bold=True)
tb(s,8.60,2.95,3.75,2.25,[("Don't resolve them — you can't.",15,False,DAMBER),("Don't drop them — they matter most.",15,False,DAMBER),("Don't assert them as verified.",15,False,DAMBER)],15,DAMBER,sp=10)
banner(s,["Named precisely, and carried untouched to the gate (Step 8) — where you answer them.","This is what stops a value judgment being smuggled in as a finding."],y=5.55,h=1.30)
note(s,"Step 3 deep-core routing.")

# ---------- S9c GRADE-LOCK + RESILIENCE (Step 4) ----------
s=slide()
label(s,"STEP 7 · GRADE-LOCK + RESILIENCE",GREEN)
headline(s,"Freeze the grades. Then try to break the call.")
card(s,0.70,1.95,5.95,3.35,fill=GTINT,border=GREEN)
tb(s,1.0,2.2,5.4,0.5,"Grade-lock",18,GREEN,bold=True)
tb(s,1.0,2.95,5.4,2.25,"Every grade from verification is now frozen. Nothing downstream can change it — a later step can only restate it. The confidence on each claim is fixed and auditable, not drifting to suit the conclusion.",16,INK)
card(s,6.85,1.95,5.78,3.35)
tb(s,7.15,2.2,5.2,0.5,"Resilience check",18,BLUE,bold=True)
tb(s,7.15,2.95,5.2,2.25,"Before committing, stress-test the call against its load-bearing assumptions and the systems map's key links: if one claim you're leaning on were wrong, would the call survive — or does everything hinge on it?",16,INK)
banner(s,"Locked grades make the confidence auditable; the resilience check makes sure the call doesn't rest on a single fragile assumption.",y=5.55,h=1.1)
note(s,"Step 4 grade-lock and resilience.")

# ---------- S10 THE FORCED GATE (dark) ----------
s=slide(NAVY_BG)
tb(s,0.72,0.4,11.9,0.4,"STEP 8 · THE FORCED GATE",15,GOLD,bold=True)
tb(s,0.70,0.78,12.0,0.9,"The value questions are yours to answer.",30,WHITE,bold=True)
card(s,0.70,1.85,11.93,1.85,fill=ATINT,border=ATINT,line=False)
tb(s,1.05,1.85,11.2,1.85,"The questions no evidence can settle — the framing, the legitimacy, the trade-offs — are handed back to you, one at a time. The analysis won't finalise until you take a position — or explicitly defer, which is logged, not ignored.",20,DAMBER,anchor=MSO_ANCHOR.MIDDLE)
tb(s,0.70,4.15,11.9,1.5,[("Why forced?  ",19,True,GOLD),
    ("Testing showed that merely showing you these questions doesn't work — attention slides to the verifiable. Being made to answer is what gets them addressed.",19,False,ICE)],19,ICE)
tb(s,0.70,5.75,11.9,0.5,"The engine surfaces the questions; you own the judgment.",17,ICE2)
note(s,"Then comes the one step the engine refuses to take from you. Running underneath every hard problem are questions no evidence can settle — how the decision is framed, who is bound by it without a say, what you are willing to trade against what. The engine separates those out, and instead of answering them, it stops and hands them to you, one at a time, and won't finalise until you take a position — or explicitly defer, which is logged, not ignored. Why force it? Because controlled testing showed that merely showing a decision-maker these questions doesn't work — they read them, then drift back to the verifiable. Being made to answer is what actually gets them addressed. The engine articulates the questions; you own the judgment.")

# ---------- S11 DECISION BRIEF ----------
s=slide()
label(s,"STEP 9 · THE DECISION BRIEF"); headline(s,"A short decision surface — not an essay.")
items=[("The call","provisional stance if the problem is wicked",CARD,BLUE),
       ("Confidence basis","how solid the EVIDENCE is — the weakest grade",GTINT,GREEN),
       ("Action contingency","whether the ACTION will work — kept separate",GTINT,GREEN),
       ("Dominant unknown","the one fact that would most move it",CARD,BLUE),
       ("Verify first","the cheapest next step that cuts risk",CARD,BLUE),
       ("What would flip it","“if this is false, the call reverses”",CARD,BLUE)]
for i,(hd,bd,fill,c) in enumerate(items):
    col=i%3; row=i//3; x=0.70+col*4.07; y=1.95+row*1.55
    card(s,x,y,3.80,1.38,fill=fill,border=GREEN if fill==GTINT else BORD)
    tb(s,x+0.24,y+0.16,3.35,0.5,hd,17,c,bold=True)
    tb(s,x+0.24,y+0.68,3.35,0.6,bd,14,BODY)
banner(s,["Every hedge names what it depends on.","A caveat that points at nothing isn't allowed."],y=5.45,h=1.30)
note(s,"With your positions on the record, it assembles the brief — not an essay, but a short brief you can act on. First the call — and where the problem is wicked, that call is a provisional stance, marked re-openable, not a verdict. Then, kept deliberately apart, two kinds of confidence people usually blur together: how solid the evidence is — naming the single weakest grade the call leans on — and whether the action will actually work, which can fail even when every fact is rock-solid. Then the one unknown that would most move the answer; the cheapest thing to check next; and what would flip the call, written as a plain condition — “I am treating this as true; if it's false, the call reverses.” Every hedge has to name what it depends on. A caveat that points at nothing isn't allowed.")

# ---------- S13b HOW THE ADJUDICATOR RUNS ----------
s=slide()
label(s,"STEP 10 · ADJUDICATION · HOW IT RUNS")
headline(s,"Run it as a subagent — blind, then adversarial.")
card(s,0.70,1.85,5.95,3.55)
tb(s,1.0,2.08,5.4,0.5,"The mechanism",18,BLUE,bold=True)
tb(s,1.0,2.75,5.4,2.6,[("• a separate subagent, its own isolated context",15,False,INK),
    ("• fed the facts + the grade-locked spine from the Anthropic analysis — a second model on top, never standalone",14,False,INK),
    ("• Pass 1 · blind — re-derive: same call?",15,False,INK),
    ("• Pass 2 · adversarial — try to break it",15,False,INK),
    ("• discrepancies fold in under grade-lock; you own the call",15,False,INK)],15,INK,sp=7)
card(s,6.85,1.85,5.78,3.55)
tb(s,7.15,2.08,5.2,0.5,"Which independent adjudicator (strongest first)",18,CORAL,bold=True)
tb(s,7.15,2.66,5.2,2.7,[
    [("A · different-LAB model (configurable — GPT-5.6 Sol, Gemini, or Grok) — strongest; runs ALONGSIDE the Anthropic model, never alone. Invoked by an API call — the engine guides the set-up interactively. ",13,False,INK),("Built + live-verified as a harness; opt-in, off by default; data leaves the boundary, so privileged/confidential is blocked by default — overridable only by a deliberate, logged act.",13,False,DAMBER,F,True)],
    ("B · a different in-house model (Fable 5)",13,False,INK),
    ("C · a panel of 2+ blind Opus instances — divergence is the signal",13,False,INK),
    ("D · a self-adversarial reset (weakest)",13,False,INK)],13,INK,sp=7)
banner(s,[("Always a check on the Anthropic model's finished analysis — the adjudicator never flips the call.  ",15,True,WHITE),("Declare the rung reached; same-model agreement is stability, not proof.",15,True,GOLD)],y=5.6,h=1.15)
note(s,"So how is that second pass actually invoked? As a subagent — a separate agent with its own isolated context, handed the facts and the grade-locked spine from the Anthropic model's analysis, but not that model's reasoning, so it can't anchor on it. Be clear on one thing: adjudication always uses two or more models. The different-lab model never runs on its own — it checks the work of the Anthropic frontier model that produced the analysis. It runs blind first — re-derive, and see if it reaches the same call — then adversarial: try to break the load-bearing claims and the call. It hands back a list of discrepancies, which fold in under grade-lock; you still own the final call. The rung it runs on depends on what is set up, strongest first. Strongest is a different-lab model — GPT-5.6 Sol, or GPT-5.5 while Sol is in limited preview — because a different training lineage gives the most decorrelated blind spots; it is built and live-verified, opt-in and off by default, and blocked by default for confidential or legal matters — overridable only by a deliberate, logged act. Next best is a different in-house model like Fable 5. If neither is set up, run a panel of two or more blind Opus instances and read where they diverge — a same-model cross-check, not independent adjudication. Weakest of all, a self-adversarial reset. Every rung is opt-in and costs credits or setup, so nothing runs on its own. Whichever rung it lands on, the engine should say so — because same-model agreement is stability, not proof.")

# ---------- S14 WHERE THE SECOND PASS EARNS ITS KEEP ----------
s=slide()
label(s,"STEP 10 · ADJUDICATION · WHERE IT HELPS"); headline(s,"Spend the frontier model where judgment decides.")
card(s,0.70,1.95,7.35,3.35,fill=GTINT,border=GREEN)
tb(s,1.0,2.2,6.8,0.5,"Sharpens the contested judgments",18,GREEN,bold=True)
tb(s,1.0,2.9,6.8,2.3,[("• load-bearing claims graded weak or uncorroborated",16,False,INK),
    ("• the disconfirmation verdicts",16,False,INK),("• any allocation of blame — where legal cases turn",16,False,INK),
    ("• the value questions at the gate",16,False,INK),("• which loop dominates, in a systems map",16,False,INK)],16,INK,sp=5)
card(s,8.30,1.95,4.33,3.35,fill=CARD)
tb(s,8.60,2.2,3.75,0.5,"Adds little on",18,BODY,bold=True)
tb(s,8.60,2.9,3.75,2.0,"the [V1] facts already confirmed against a primary source — those are settled.",17,BODY)
banner(s,["It reliably sharpens the contested judgments — safely, without reversing a sound call.","Rank its objections to the few that change the decision."],y=5.50,h=1.30)
note(s,"That second pass doesn't redo the whole analysis — it goes where judgment, not fact, decides the answer: the claims graded weak or uncorroborated, the disconfirmation verdicts, any allocation of blame, the value questions at the gate, which loop dominates. There it reliably sharpens — catching where a dominance claim isn't yet established, where an exonerating metric is itself confounded, where an argument contradicts itself — and it does so without reversing a sound call. One caveat from testing: it is verbose, so rank its objections and surface only the few that would actually change the decision.")

# ---------- S13.5 YOU CONTROL IT (governance + retirement) ----------
s=slide()
label(s,"STEP 10 · ADJUDICATION · YOU CONTROL IT")
headline(s,"Every gate is yours to set — and it retires itself.")
card(s,0.70,1.95,5.95,3.35)
tb(s,1.0,2.2,5.4,0.5,"Switchable — nothing is locked",18,BLUE,bold=True)
tb(s,1.0,2.9,5.4,2.3,[("• offer or not; which rung; panel size N, labs M, runs R",14,False,INK),
    ("• the privilege / egress gate is asked before anything leaves the boundary",14,False,INK),
    ("• a block lifts only by a deliberate, logged act — never silently",14,False,INK),
    ("• the exact redacted package is shown to confirm / edit / cancel",14,False,INK),
    ("• every setting standing or per-run; personal-use, nothing centrally locked",14,False,INK)],14,INK,sp=6)
card(s,6.85,1.95,5.78,3.35,fill=GTINT,border=GREEN)
tb(s,7.15,2.2,5.2,0.5,"It earns its place — or it retires",18,GREEN,bold=True)
tb(s,7.15,2.9,5.2,2.3,[("• every real run is logged to a use-monitor",14,False,INK),
    ("• a pre-registered, per-class rule reads the ledger",14,False,INK),
    ("• no decision-relevant catches for a class over time — recommend on-request (reversible, audited)",14,False,INK),
    ("• decides economic redundancy from real use, not by argument",14,False,INK)],14,INK,sp=6)
banner(s,[("You own even the boundary decision — ",15,True,WHITE),("and the layer steps back where it no longer adds decision value.",15,True,GOLD)],y=5.6,h=1.15)
note(s,"One more thing about that second pass before we step back to the model itself: you are in control of all of it, and it is honest about when it is not needed. Every part is switchable, and nothing is locked — whether it is offered at all, which rung runs, how many Opus instances, how many different labs, how many runs each. All yours to set, standing or per run. The governance is asked-first: before anything leaves Anthropic's boundary, the engine asks whether the matter is privileged or confidential, and shows you the exact redacted package to confirm, edit, or cancel. A privileged block is the default, and it lifts only by a deliberate, logged act — never silently. This is a personal-use tool: you own even the boundary decision. And the layer is honest about its own worth: every real run is logged to a use-monitor, and a pre-registered, per-class rule reads that ledger — if a whole class of problem shows no decision-relevant catches over enough runs, the rule recommends demoting adjudication to on-request for that class — reversible, and still spot-audited — deciding from real use rather than argument where it no longer adds decision value.")

# ---------- S13 WHICH MODEL, AND WHEN ----------
s=slide()
label(s,"MODEL CHOICE · WHICH MODEL, AND WHEN"); headline(s,"Which model, and when.")
card(s,0.70,1.95,5.95,3.05)
tb(s,1.0,2.2,5.4,0.9,"Execution — a top-tier frontier model",19,BLUE,bold=True)
tb(s,1.0,3.55,5.4,1.4,[("Nearly every step leans on the model's own reasoning.",16,False,INK),
    ("We ran the engine on Opus 4.8.",16,True,BLUE)],16,INK,sp=6)
card(s,6.85,1.95,5.78,3.05,fill=GTINT,border=GREEN)
tb(s,7.15,2.2,5.2,0.9,"High-stakes adjudication — a second, independent frontier model",19,GREEN,bold=True)
tb(s,7.15,3.55,5.2,1.4,[("For legal exposure, contested blame, decisions that bind people.",16,False,INK),
    ("We used Fable 5 — run blind, then adversarial.",16,True,GREEN)],16,INK,sp=6)
banner(s,"A rigour-and-defensibility check, not a safety net: it sharpens a sound analysis, safely — but catching what the disciplined first pass missed is unproven.",y=5.2,h=1.35)
note(s,"One practical question: which model should run this? The engine is designed for a strong model — a top-tier frontier model for execution, because nearly every step leans on the model's own reasoning: judging a source, running a real search for counter-evidence, catching its own confident mistakes. In building the engine we ran it on Opus 4.8. But for the highest-stakes work — legal exposure, contested blame, a decision that binds people — add a second pass: a different frontier model, brought in to adjudicate. We used Fable 5. Here's the honest, tested part: its benefit at catching an error the disciplined first pass missed is unproven — on our own testing, the first pass already caught the constructible ones. What the second pass does do, reliably and safely, is sharpen a sound analysis: it flags where a load-bearing claim isn't yet established, without corrupting a correct call. Treat it as a rigour-and-defensibility check, not a safety net.")

# ---------- S15 ON A WEAKER MODEL ----------
s=slide()
label(s,"MODEL CHOICE · A WEAKER MODEL",AMBER)
headline(s,"The engine disciplines the model — it doesn't rescue it.")
card(s,0.70,1.95,5.95,2.6,fill=GTINT,border=GREEN)
tb(s,1.0,2.18,5.4,0.5,"Survives — the scaffolding",18,GREEN,bold=True)
tb(s,1.0,2.85,5.4,1.6,[("• the forced gate",16,False,INK),("• grades & caveats still travel",16,False,INK),
    ("• audience rendering",16,False,INK)],16,INK,sp=5)
card(s,6.85,1.95,5.78,2.6,fill=ATINT,border=CORAL)
tb(s,7.15,2.18,5.2,0.5,"Degrades — what fills it",18,AMBER,bold=True)
tb(s,7.15,2.85,5.2,1.6,[("• source tiers less reliable",16,False,INK),("• thinner counter-evidence search",16,False,INK),
    ("• more confident errors, worse at catching them",16,False,INK)],16,INK,sp=5)
banner(s,[("The polish is itself a risk — grades and caveats can make a weak answer look more defensible than it is. ",17,True,WHITE),("The assurance is only as good as the model underneath.",17,True,GOLD)],y=4.85,h=1.6)
note(s,"And the other direction — what if you run it on a weaker model? The scaffolding survives: the gate still makes you own the value calls, the grades and caveats still travel. But what fills that scaffold degrades. A weaker model grades sources less reliably, runs a thinner search for counter-evidence, and carries more confident errors while being worse at catching them — and the engine relies on the model catching itself. The warning worth stating plainly: the engine cannot add reasoning the model doesn't have. It disciplines and exposes the model's work; it does not rescue it. And its own polish is a risk — the grades and caveats can make a weak answer look more defensible than it is. The assurance is only ever as good as the model underneath it.")

# ---------- S12 FIVE COMMANDS ----------
s=slide()
label(s,"THE INTERFACE"); headline(s,"Five commands.")
cmds=[("/insight-engine:analyse","Runs the whole pipeline, end to end."),
      ("/insight-engine:verify","Fact-checks a set of claims, or a single premise."),
      ("/insight-engine:render","Re-voices a finished analysis for a reader — never changing a grade."),
      ("/insight-engine:track","Keeps an investigation live as new facts arrive."),
      ("/insight-engine:adjudicate","Opt-in, off-by-default independent second pass (Step 10) on a finished high-stakes analysis.")]
for i,(cm,bd) in enumerate(cmds):
    col=i%2; row=i//2; x=0.70+col*6.07; y=1.78+row*1.72
    card(s,x,y,5.80,1.55)
    tb(s,x+0.32,y+0.22,5.25,0.5,cm,17,BLUE,bold=True,font="Consolas")
    tb(s,x+0.32,y+0.80,5.2,0.6,bd,14,BODY)
note(s,"All of this runs through five commands. Analyse runs the whole pipeline you've just seen. Verify runs the fact-checking layer on its own — a set of claims, or a single premise you want pressure-tested. Render re-voices a finished analysis for a particular reader — a board, a regulator, opposing counsel, a family — without ever changing a claim or a grade to suit the audience, and carrying a core of caveats no version is allowed to drop. Track keeps an investigation alive as new facts arrive, re-verifying the spine and telling you if the call has moved. And adjudicate is the opt-in independent second pass — the L4.5 layer — that brings in one or more separate models to check a finished high-stakes analysis.")

# ---------- S14b RENDER (Step 11) ----------
s=slide()
label(s,"STEP 11 · RENDER")
headline(s,"After the call: re-voice it for the reader.")
card(s,0.70,1.95,11.93,3.25,fill=CARD,border=BLUE)
tb(s,1.0,2.2,11.3,0.6,"Render — for a particular reader",19,BLUE,bold=True)
tb(s,1.0,3.0,11.3,2.1,"Re-voice the finished analysis for whoever must read it — a board, a regulator, opposing counsel, a family — without ever changing a claim or a grade to suit the audience, and carrying a core of caveats that no version is allowed to drop.",16,BODY)
banner(s,"A register transform only — it changes the voice, never a claim or a grade.",y=5.55,h=1.1)
note(s,"S14b render note")

# ---------- S14c TRACK (Step 12) ----------
s=slide()
label(s,"STEP 12 · TRACK",GREEN)
headline(s,"Keep the analysis alive as facts arrive.")
card(s,0.70,1.95,11.93,3.25,fill=GTINT,border=GREEN)
tb(s,1.0,2.2,11.3,0.6,"Track — as new facts arrive",19,GREEN,bold=True)
tb(s,1.0,3.0,11.3,2.1,"Keep a long-running investigation alive: re-verify the grade-locked spine as fresh evidence lands, and say plainly if the call has moved. So the analysis does not simply end at the brief — it is kept current as the world changes.",16,BODY)
banner(s,"The graded spine is re-checked on every update — a living analysis, not a snapshot.",y=5.55,h=1.1)
note(s,"S14c track note")

# ---------- S17 POSITIONING ----------
s=slide()
label(s,"POSITIONING · ICD 203",GREEN); headline(s,"Architected around ICD 203.")
tb(s,0.70,2.15,6.9,2.6,"ICD 203 is the U.S. intelligence community's Analytic Standards directive. The engine is built around its nine analytic tradecraft standards — describing source quality, expressing uncertainty, separating evidence from judgment, weighing alternatives.",19,INK)
tb(s,0.70,4.35,6.9,1.5,"Architected around the standards and assessed once — not a formal certification. A signal, not a seal.",16,BODY,italic=True)
card(s,8.10,1.9,4.53,4.4,fill=CARD)
tb(s,8.35,2.2,4.05,0.5,"BLIND ASSESSMENT · ICD 203",14,BODY,bold=True)
tb(s,8.35,2.65,4.05,1.3,"92%",72,GREEN,bold=True)
tb(s,8.35,3.95,4.05,0.4,"the Insight Engine",16,INK)
tb(s,8.35,4.7,4.05,0.9,"75%",40,BODY,bold=True)
tb(s,8.35,5.55,4.05,0.4,"a strong model unaided",16,BODY)
note(s,"None of this is invented from nothing. The engine is architected around ICD 203 — the U.S. intelligence community's Analytic Standards directive — and specifically its nine analytic tradecraft standards: describing the quality of sources, expressing uncertainty, separating evidence from judgment, weighing alternatives. In a blind assessment against those nine standards, the engine scored ninety-two percent, where a strong model unaided scored seventy-five. To be honest about what that is: it's architected around ICD 203 and assessed once — not a formal certification. The number is a signal, not a seal.")

# ---------- S18 CLOSE (dark) ----------
s=slide(NAVY_BG)
tb(s,0.90,2.1,11.5,1.1,"Bring it your hardest problem.",48,WHITE,bold=True)
tb(s,0.92,3.45,11.0,1.3,"An open plugin for Claude Code & Cowork. Hand it a real problem, let the gate make you take a position — and get back an analysis you can defend.",20,ICE)
tb(s,0.92,5.15,8.0,0.6,"github.com/ComplexAR/Insight-Engine",18,WHITE,bold=True,font="Consolas")
tb(s,0.92,6.5,11.4,0.5,"The large language model provides the insight; the engine makes it verified, defensible, and audience-ready.",17,ICE2)
note(s,"So — bring it your hardest problem. It's an open plug-in for Claude Code and Cowork. Hand it a real problem, let the gate make you take a position, and get back an analysis you can stand behind.")

out=os.path.join(os.path.dirname(os.path.abspath(__file__)),"Insight-Engine-Explainer.pptx")
_tot=len(prs.slides._sldIdLst)
for _i,_s in enumerate(prs.slides,1):
    _dark=False
    try: _dark=(str(_s.background.fill.fore_color.rgb)=="141413")
    except: pass
    tb(_s,11.83,7.02,1.1,0.33,f"{_i}/{_tot}",11,(ICE2 if _dark else MUTE),align=PP_ALIGN.RIGHT)
import json as _json
_json.dump(CARDGEO, open(os.path.join(os.path.dirname(os.path.abspath(__file__)),"card_geometry.json"),"w"))
prs.save(out)
print("Saved", out, "slides:", len(prs.slides._sldIdLst))
